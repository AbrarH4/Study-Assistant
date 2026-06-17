import os
import sys
import json
import re
import urllib.request
from pathlib import Path
import threading
from dotenv import load_dotenv
from tkinter import filedialog, messagebox
import customtkinter as ctk
import hashlib
import pickle
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
# =====================================================================
# PREMIUM THEME CONFIGURATION (Silicon Valley Corporate Aesthetic)
# =====================================================================
ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")
# =====================================================================
load_dotenv()
PROVIDERS = [
    {
        "name": "Google AI Studio",
        "url": "https://generativelanguage.googleapis.com/v1beta/openai/",
        "key": os.getenv("GEMINI_API_KEY"),
        "model": "gemini-2.5-flash",
    },
    {
        "name": "Groq Cloud",
        "url": "https://api.groq.com/openai/v1",
        "key": os.getenv("GROQ_API_KEY"),
        "model": "llama-3.3-70b-versatile",
    },
    {
        "name": "OpenRouter Free",
        "url": "https://openrouter.ai/api/v1",
        "key": os.getenv("OPENROUTER_API_KEY"),
        "model": "openrouter/free",
    },
]
# =====================================================================

# Premium User-Defined Palette Hex Codes
BG_MAIN = "#0D0E12"        # Deep Obsidian Base
BG_CARD = "#161920"        # Elevated Slate Containers
BORDER_COLOR = "#262930"   # Subtle crisp border
ACCENT_COLOR = "#4F46E5"   # Premium Indigo Action Button
ACCENT_HOVER = "#4338CA"   # Deep Indigo Hover

TEXT_PRIMARY = "#E2E8F0"   # Crisp White-Slate (High Contrast Headers)
TEXT_SECONDARY = "#94A3B8" # Muted Slate Gray (Comfortable Body Text)
TEXT_DISABLED = "#475569"  # Receding Hint/Placeholder Gray
TEXT_INTERACTIVE = "#F2F0E8" # Saffron Yellow Accent Highlight

# =====================================================================
# LOGIC ENGINE
# =====================================================================
BASE_DIR = Path(__file__).parent

SETTINGS_FILE = BASE_DIR / "settings.json"
STOP_WORDS_FILE = BASE_DIR / "STOP_WORDS.json"

APP_VERSION = "1.0.0"
EMBEDDINGS_CACHE_FILE = BASE_DIR / "embeddings_cache.pkl"
#  ----------------------------------
Notes = {}


model = None
model_loaded = False
is_online = False

def get_hash(filepath):
    with open(filepath,'rb') as f:
        return hashlib.md5(f.read()).hexdigest()
def check_internet_connection(timeout=3):
    try:
        urllib.request.urlopen("https://google.com", timeout=timeout)
        return True
    except Exception:
        return False


def bg_model_loading(ui_instance):
    """Thread worker that loads models and checks workspace safely without freezing UI."""
    global model, model_loaded, is_online
    is_online = check_internet_connection()

    if not is_online:
        ui_instance.after(
            0,
            lambda: ui_instance.display_critical_network_error(
                "CRITICAL: NO INTERNET CONNECTION DETECTED. Please check your network and restart the application."
            ),
        )
        return

    ui_instance.after(
        0, lambda: ui_instance.update_loading_status("LOADING EMBEDDINGS ENGINE...")
    )

    from sentence_transformers import SentenceTransformer

    model = SentenceTransformer("all-MiniLM-L6-v2")

    if current_notes_path:
        load_notes_from_path(current_notes_path)

    model_loaded = True
    ui_instance.after(0, ui_instance.transition_to_main_ui)
# PDF TEXT EXTRACTION ------------------------------------------>
def pdf_loader(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        pages = pdf.pages
        for page in pages:
            text += page.extract_text() or ""
    return text.strip()
# DOCX FILE TEXT EXTRACTION -------------------------------->
def docx_loader(filepath):
    text =""
    doc  = Document(filepath)
    for paragraphs in doc.paragraphs:
        text += paragraphs.text or ""
    return text.strip()
# PRESENTATION FILE TEXT EXTRACTION -------------------------------->
def presentation_loader(file_path):
    text = ""
    prs = Presentation(file_path)
    for slides in prs.slides:
        for shapes in slides.shapes:
            if shapes.has_text_frame:
                for paragraph in shapes.text_frame.paragraphs:
                    text += paragraph.text or ''
    return text.strip()
# LOADING NOTES FOLDER ------------------->        
def load_folder():
    if SETTINGS_FILE.exists():
        with open(SETTINGS_FILE, "r", encoding="utf-8") as f:
            settings = json.load(f)
            if "path_notes" in settings and settings["path_notes"]:
                if os.path.exists(settings["path_notes"]):
                    return settings.get("path_notes")
                    
            
    path_notes = filedialog.askdirectory(title="CHOOSE YOUR NOTES FOLDER.")
    if path_notes:
        with open(SETTINGS_FILE, "w", encoding="utf-8") as f:
            json.dump({"path_notes": path_notes}, f, indent=4)
    return path_notes


Embedding_cache = {}
Notes_Cache = {}


def load_notes_from_path(folder_path):
    """Helper method to completely re-index target files into memory."""
    ALLOWED_EXTENSIONS = {
        ".txt",
        ".md",
        '.pdf',
        '.docx',
        '.pptx'
    }
    global Notes
    Notes.clear()
    Embedding_cache.clear()
    Notes_Cache.clear()
    disk_cache = {}
    if EMBEDDINGS_CACHE_FILE.exists():
        with open(EMBEDDINGS_CACHE_FILE,'rb') as f:
            disk_cache = pickle.load(f)
    notes_changed = False
    for files in os.listdir(folder_path):
        path = os.path.join(
            folder_path,
            files
        )
        
        if (Path(path).suffix.lower() not in ALLOWED_EXTENSIONS):
            continue
        if os.path.isfile(path):
            try:
                if Path(path).suffix.lower() == ".pdf":
                    content = pdf_loader(path)
                elif Path(path).suffix.lower() == ".docx":
                    content = docx_loader(path)
                elif Path(path).suffix.lower() == ".pptx":
                    content = presentation_loader(path)
                else:    
                    with open(path, "r", encoding="utf-8") as f:
                        content = f.read()
                Notes[files] = content
                file_hash = get_hash(path)
                if files in disk_cache and disk_cache[files][0] == file_hash:
                    Embedding_cache[files] = disk_cache[files][1]
                else:
                    Embedding_cache[files] = model.encode(content,convert_to_tensor=True)
                    disk_cache[files] = (file_hash,Embedding_cache[files])
                    notes_changed = True  
            except Exception as e:
                print(f"Skipping unreadable asset {files}: {e}")
    disk_cache = {k:v for k,v in disk_cache.items() if k in Notes}
    if notes_changed:
        with open(EMBEDDINGS_CACHE_FILE,'wb') as f:
            pickle.dump(disk_cache,f)

current_notes_path = load_folder()
if current_notes_path:

    print(
        f"Loaded notes folder: "
        f"{Path(current_notes_path).name}"
    )

else:

    print(
        "No notes folder configured."
    )
    messagebox.showinfo(
    "One Quick Step",
    "Study Assistant needs to know where your notes live.\n\nHit ⚙ Path to point it to your notes folder — you only need to do this once."
)
    


def load_stop_words():
    try:
        with open(
            STOP_WORDS_FILE,
            "r",
            encoding="utf-8"
        ) as file:

            word_list = json.load(file)
            return set(word_list)

    except FileNotFoundError:
        print(
            "WARNING: STOP_WORDS.json not found."
        )
        return set()

    except Exception as e:
        print(
            f"Failed to load stop words: {e}"
        )
        return set()


Stop_words = load_stop_words()


def keyword(question: str = None) -> list:
    if question is None:
        question = input()
    words = question.split()
    keywords = [
        word.strip("?!.,").lower()
        for word in words
        if word.lower() not in Stop_words
    ]
    return keywords


def Ranking_System(keyword, question):
    if not keyword:
        messagebox.showerror("ERROR", "SORRY WE COULD'NT FIND THE BEST NOTES FOR YOU.")
        return None

    from sentence_transformers import util

    scores = {}
    semantic_scores = {}

    for keys, notes in Notes.items():
        if keys.lower() == "error.txt":
            continue

        scores[keys] = 0
        notes_lower = notes.lower()
        for word in keyword:
            word_lower = word.lower()
            if word_lower in keys.lower():
                scores[keys] += 20
            if word_lower in notes_lower:
                scores[keys] += 1

       

    encoded_question = model.encode(question, convert_to_tensor=True)

    for keys, notes in Notes.items():
        if keys.lower() == "error.txt":
            continue
        encoded_notes = Embedding_cache[keys]
        score = util.cos_sim(encoded_question, encoded_notes)
        semantic_scores[keys] = score.item()

    Final_Scores = {}
    for key in scores:
        if key.lower() == "error.txt":
            continue
        combined_score = scores[key] + semantic_scores.get(key, 0.0)
        Final_Scores[key] = combined_score

    if not Final_Scores:
        # FIX #1: return a consistent (list_or_str, tensor) tuple so the caller's
        # unpacking `winning_file_name, encoded_question_tensor = ranking_result`
        # never crashes.
        return "error.txt", encoded_question

    best_note = [
        note_name
        for note_name, score in sorted(
            Final_Scores.items(), key=lambda x: x[1], reverse=True
        )[:3]
    ]

    print(f"Selected Notes: {best_note}")
    return best_note, encoded_question


from openai import OpenAI


def GenerateAnswer(question, context):
    system_prompt = (
        "You are an intelligent and highly capable study assistant.\n"
        "Your primary source of truth is the provided notebook context.\n\n"
        "CORE RULES:\n"
        "1. Use the notebook context as the factual source for your answers.\n"
        "2. You may rephrase, simplify, summarize, and explain information in your own words.\n"
        "3. You may adapt explanations to different learning levels and teaching styles.\n"
        "4. Do NOT invent facts that are not supported by the notebook context.\n"
        "5. Do NOT contradict the notebook context.\n"
        "6. If the notebook context contains partial information, answer using the available information and clearly mention any missing details.\n"
        "7. Only say information is unavailable when the notebook context contains no relevant information at all.\n"
        "8. Prioritize understanding over memorization.\n"
        "9. Use examples whenever they improve understanding.\n"
        "10. Keep answers clear, accurate, and educational.\n\n"
        "TEACHING MODES:\n"
        "1. If the user asks for a simple explanation, beginner explanation, ELI5 explanation, or asks to explain something like they are a child, use simple language, analogies, and everyday examples.\n"
        "2. If the user asks for a detailed, technical, advanced, or university-level explanation, provide deeper detail using the notebook context.\n"
        "3. If the user asks for examples, provide examples whenever possible.\n"
        "4. If the user asks for differences or comparisons, compare the concepts using only information supported by the notebook context.\n\n"
        "ANSWERING STYLE:\n"
        "1. Start with a direct answer whenever possible.\n"
        "2. Follow with a short explanation.\n"
        "3. Add examples when helpful.\n"
        "4. Avoid copying the notes word-for-word unless necessary.\n"
        "5. Focus on helping the user understand the concept.\n"
        "6. Use clean formatting with paragraphs and bullet points when useful.\n\n"
        "CRITICAL PLAIN-TEXT BOX LAYOUT OVERRIDES (MANDATORY):\n"
        "To ensure your text formats correctly in our custom UI layout engine, you must use these explicit markers:\n"
        "- Every primary concept title, heading, or major section MUST begin exactly with '## ' (e.g., ## Softmax vs Argmax).\n"
        "- Before starting any major topic change or comparison section, place a single line containing only '---' to split the layout visually.\n"
        "- For point-by-point lists or bullet layouts, use a clean dash followed by a space: '- Your point here'.\n"
        "- Do NOT wrap lines in thick bold markdown '**' inside headers or bullets.\n\n"
        "IMPORTANT:\n"
        "The notebook context contains the knowledge.\n"
        "Your job is to teach that knowledge clearly.\n"
        "Do not use outside knowledge.\n"
        "Do not fabricate information.\n"
        "If relevant information exists in the context, explain it naturally instead of refusing to answer."
    )

    user_content = f"QUESTION:\n{question}\n\nNOTEBOOK CONTEXT:\n{context}"

    for provider in PROVIDERS:
        if not provider["key"]:
            print(
                f"Skipping {provider['name']}: No API key detected in your .env file."
            )
            continue

        try:
            print(f"Connecting to AI pipeline: {provider['name']}...")
            client = OpenAI(base_url=provider["url"], api_key=provider["key"])
            response = client.chat.completions.create(
                model=provider["model"],
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content},
                ],
                temperature=0.1,
            )
            print(f"[SUCCESS] Generated answer via {provider['name']}!")
            return response.choices[0].message.content

        except Exception as error:
            print(f"!!! Warning: {provider['name']} failed or was rate-limited: {error}")
            print("Cascading down to the next available provider...")
            continue

    return "CRITICAL FAILURE: All configured online AI service pipelines are currently exhausted or unreachable."


def AnswerSystem(winning_file_names, encoded_question, original_question):
    try:
        all_context = ""
        for filename in winning_file_names:
            if filename in Notes:
                all_context += f"SOURCE: {filename}\n{Notes[filename]}\n\n"

        if not all_context:
            return "No relevant context found in selected notes."

        ai_response = GenerateAnswer(original_question, all_context)
        return ai_response if ai_response else all_context

    except Exception as e:
        print(f"Extraction Error: {e}")
        return f"An operational breakdown occurred: {str(e)}"
# =====================================================================
# UI ARCHITECTURE
# =====================================================================


class StudyAssistantUI(ctk.CTk):

    def __init__(self):
        super().__init__()
        self.title("Study Assistant")
        self.after(0, lambda: self.wm_state("zoomed"))
        if sys.platform.startswith("linux"):
            self.attributes("-zoomed", True)
        else:
            self.state("zoomed")

        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(0, weight=1)

        # ------------------- SPLASH SCREEN -------------------
        self.loading_frame = ctk.CTkFrame(self, fg_color=BG_MAIN)
        self.loading_frame.grid(row=0, column=0, sticky="nsew")
        self.loading_frame.grid_rowconfigure(0, weight=1)
        self.loading_frame.grid_rowconfigure(1, weight=0)
        self.loading_frame.grid_rowconfigure(2, weight=0)
        self.loading_frame.grid_rowconfigure(3, weight=0)
        self.loading_frame.grid_rowconfigure(4, weight=0)
        self.loading_frame.grid_rowconfigure(5, weight=1)
        self.loading_frame.grid_columnconfigure(0, weight=1)

        # App name — large, centered, indigo-tinted white
        self.loading_app_name = ctk.CTkLabel(
            self.loading_frame,
            text="Study Assistant",
            font=ctk.CTkFont(family="Segoe UI", size=42, weight="bold"),
            text_color="#E8EEFF",
        )
        self.loading_app_name.grid(row=1, column=0, pady=(0, 6))

        # Tagline — small, muted, below the name
        self.loading_tagline = ctk.CTkLabel(
            self.loading_frame,
            text="Your personal AI-powered notebook",
            font=ctk.CTkFont(family="Segoe UI", size=14),
            text_color="#4F6080",
        )
        self.loading_tagline.grid(row=2, column=0, pady=(0, 48))

        # Status text — updates as each stage completes
        self.loading_text = ctk.CTkLabel(
            self.loading_frame,
            text="Checking connection…",
            font=ctk.CTkFont(family="Segoe UI", size=12),
            text_color="#475569",
        )
        self.loading_text.grid(row=3, column=0, pady=(0, 10))

        # Thin progress bar — full width, barely 4px tall
        self.progress_bar = ctk.CTkProgressBar(
            self.loading_frame,
            height=3,
            corner_radius=2,
            fg_color="#1E2433",
            progress_color=ACCENT_COLOR,
        )
        self.progress_bar.grid(row=4, column=0, padx=200, pady=(0, 0), sticky="ew")
        self.progress_bar.configure(mode="indeterminate")
        self.progress_bar.start()

        # Version tag — bottom left, barely visible
        self.loading_version = ctk.CTkLabel(
            self.loading_frame,
            text=f"v{APP_VERSION}",
            font=ctk.CTkFont(family="Segoe UI", size=11),
            text_color="#262930",
        )
        self.loading_version.grid(row=5, column=0, pady=(0, 20))
        # ------------------- MAIN INTERFACE (Hidden Initially) -------------------
        self.main_canvas = ctk.CTkFrame(self, fg_color="transparent")
        self.main_canvas.grid_columnconfigure(0, weight=1)
        self.main_canvas.grid_rowconfigure(3, weight=1)

        self.header_label = ctk.CTkLabel(
            self.main_canvas,
            text=f"Study Assistant v{APP_VERSION}",
            font=ctk.CTkFont(family="Segoe UI", size=28, weight="bold"),
            text_color=TEXT_PRIMARY,
            anchor="center",
        )
        self.header_label.grid(row=0, column=0, pady=(0, 5), sticky="ew")

        display_path = (
            current_notes_path if current_notes_path else "No Path Indexed"
        )
        self.sub_header = ctk.CTkLabel(
            self.main_canvas,
            text=f"• DIRECTORY: {Path(display_path).name}",
            font=ctk.CTkFont(family="Segoe UI", size=12),
            text_color="#677EB6",
            anchor="center",
        )
        self.sub_header.grid(row=1, column=0, pady=(0, 20), sticky="ew")

        self.control_panel = ctk.CTkFrame(
            self.main_canvas,
            fg_color=BG_CARD,
            border_color=BORDER_COLOR,
            border_width=1,
            corner_radius=12,
        )
        self.control_panel.grid(row=2, column=0, pady=(0, 20), sticky="ew")
        self.control_panel.grid_columnconfigure(0, weight=1)

        self.search_entry = ctk.CTkEntry(
            self.control_panel,
            placeholder_text="Ask your question here...",
            height=45,
            font=ctk.CTkFont(family="Segoe UI", size=13),
            fg_color=BG_MAIN,
            border_color=BORDER_COLOR,
            text_color=TEXT_PRIMARY,
            placeholder_text_color=TEXT_DISABLED,
            corner_radius=8,
            justify="center",
        )
        self.search_entry.grid(
            row=0, column=0, padx=(20, 10), pady=20, sticky="ew"
        )
        self.search_entry.bind(
            "<Return>", lambda event: self.ui_trigger_search_flow()
        )

        self.btn_frame = ctk.CTkFrame(
            self.control_panel, fg_color="transparent"
        )
        self.btn_frame.grid(row=0, column=1, padx=(0, 20), pady=20)

        self.search_btn = ctk.CTkButton(
            self.btn_frame,
            text="Analyze",
            height=45,
            width=100,
            font=ctk.CTkFont(family="Segoe UI", size=14, weight="bold"),
            fg_color=ACCENT_COLOR,
            hover_color=ACCENT_HOVER,
            text_color="#FFFFFF",
            corner_radius=8,
            command=self.ui_trigger_search_flow,
        )
        self.search_btn.grid(row=0, column=0, padx=(0, 5))

        self.path_btn = ctk.CTkButton(
            self.btn_frame,
            text="⚙ Path",
            height=45,
            width=80,
            font=ctk.CTkFont(family="Segoe UI", size=13, weight="bold"),
            fg_color="#1F2937",
            hover_color="#374151",
            text_color=TEXT_PRIMARY,
            corner_radius=8,
            command=self.ui_change_path_flow,
        )
        self.path_btn.grid(row=0, column=1)

        # ── OUTPUT PANEL ──────────────────────────────────────────────
        self.status_box = ctk.CTkTextbox(
            self.main_canvas,
            activate_scrollbars=True,
            fg_color="#0F1117",          # slightly cooler than BG_CARD — gives the
            border_color="#1E2433",      # output pane its own identity
            border_width=1,
            text_color="#C8D0E0",        # softer than TEXT_PRIMARY — easier on long reads
            font=ctk.CTkFont(family="Segoe UI", size=14),
            corner_radius=14,
            spacing1=4,
            spacing2=2,
            spacing3=8,
        )
        self.status_box.grid(row=3, column=0, sticky="nsew")

        tb = self.status_box._textbox   # shorthand for all tag_config calls below

        # ── TYPOGRAPHIC TAG SYSTEM ─────────────────────────────────────
        # H1 — big concept title, indigo-tinted white, generous leading
        tb.tag_config(
            "h1",
            font=("Segoe UI", 22, "bold"),
            foreground="#E8EEFF",
            spacing1=18,
            spacing3=4,
        )
        # H2 — section subheading, muted indigo accent
        tb.tag_config(
            "h2",
            font=("Segoe UI", 15, "bold"),
            foreground="#818CF8",        # indigo-400 — connects to ACCENT_COLOR family
            spacing1=14,
            spacing3=2,
        )
        # Body — comfortable reading gray, slightly warm
        tb.tag_config(
            "body",
            font=("Segoe UI", 14),
            foreground="#C8D0E0",
            spacing1=2,
            spacing3=4,
        )
        # Bullet — soft cyan lead character + white text
        tb.tag_config(
            "bullet_marker",
            font=("Segoe UI", 14, "bold"),
            foreground="#38BDF8",        # sky-400 — cool contrast against indigo headings
        )
        tb.tag_config(
            "bullet_text",
            font=("Segoe UI", 14),
            foreground="#CBD5E1",
            spacing3=3,
        )
        # Divider — barely visible rule, purely structural
        tb.tag_config(
            "divider",
            font=("Segoe UI", 6),
            foreground="#1E2433",
            spacing1=10,
            spacing3=10,
        )
        # Loading pulse
        tb.tag_config(
            "loading",
            font=("Segoe UI", 13, "italic"),
            foreground="#818CF8",
        )
        # Source footer — smallest, most receded
        tb.tag_config(
            "source_label",
            font=("Segoe UI", 11, "bold"),
            foreground="#475569",
            spacing1=14,
            spacing3=2,
        )
        tb.tag_config(
            "source_item",
            font=("Segoe UI", 12),
            foreground="#4F6080",
            spacing3=2,
        )

        self.status_box.insert("1.0", "Ready.")
        self.status_box.configure(state="disabled")

        threading.Thread(target=bg_model_loading, args=(self,), daemon=True).start()

    def display_critical_network_error(self, message):
        self.progress_bar.stop()
        self.progress_bar.configure(progress_color="#EF4444")
        self.loading_text.configure(text=message, text_color="#EF4444")
        self.loading_tagline.configure(
            text="Check your network and restart.", text_color="#475569"
        )

    def update_loading_status(self, text_status):
        self.loading_text.configure(text=text_status)

    def transition_to_main_ui(self):
        self.progress_bar.stop()
        self.loading_frame.grid_forget()
        self.main_canvas.grid(row=0, column=0, padx=40, pady=30, sticky="nsew")

    def ui_change_path_flow(self):
        path_notes = filedialog.askdirectory(title="CHOOSE YOUR NEW NOTES FOLDER.")
        if not path_notes:
            return
        if not model_loaded:
            messagebox.showerror("ERROR", "Model is still loading, please wait.")
            return

        with open(SETTINGS_FILE, "w", encoding="utf-8") as f:
            json.dump({"path_notes": path_notes}, f, indent=4)

        self.status_box.configure(state="normal")
        self.status_box.delete("1.0", ctk.END)
        self.status_box._textbox.insert("1.0", f"Re-indexing  {Path(path_notes).name}…", "loading")
        self.status_box.configure(state="disabled")

        self.path_btn.configure(state="disabled")
        self.search_btn.configure(state="disabled")

        def reindex_worker():
            load_notes_from_path(path_notes)
            self.after(0, lambda: self.finish_reindex(path_notes))

        threading.Thread(target=reindex_worker, daemon=True).start()

    def finish_reindex(self, path_notes):
        self.sub_header.configure(text=f"• DIRECTORY: {Path(path_notes).name}")
        self.status_box.configure(state="normal")
        self.status_box.delete("1.0", ctk.END)
        self.status_box._textbox.insert("1.0", f"Indexed  →  {path_notes}", "body")
        self.status_box.configure(state="disabled")
        self.path_btn.configure(state="normal")
        self.search_btn.configure(state="normal")

    def ui_trigger_search_flow(self):
        if not model_loaded:
            return
        query_text = self.search_entry.get()
        self.search_entry.delete(0, ctk.END)
        if not query_text.strip():
            return

        self.search_btn.configure(state="disabled")
        self.search_entry.configure(state="disabled")
        self.status_box.configure(state="normal")
        self.status_box.delete("1.0", ctk.END)
        self.status_box.configure(state="disabled")

        is_loading = True

        def animate_terminal_loading(tick=0):
            if not is_loading:
                return
            frames = ["⠋", "⠙", "⠸", "⠴", "⠦", "⠇"]
            spinner = frames[tick % len(frames)]
            self.status_box.configure(state="normal")
            self.status_box.delete("1.0", ctk.END)
            self.status_box._textbox.insert(
                "1.0", f"  {spinner}  Thinking…", "loading"
            )
            self.status_box.configure(state="disabled")
            self.after(120, lambda: animate_terminal_loading(tick + 1))

        animate_terminal_loading()

        def async_search_pipeline():
            nonlocal is_loading
            global input
            old_input = input
            input = lambda: query_text

            extracted_keywords = keyword(query_text)
            ranking_result = Ranking_System(extracted_keywords, query_text)
            input = old_input

            if not ranking_result:
                is_loading = False
                self.status_box.after(
                    0, lambda: self.render_fallback_msg(
                        "No matching notes found for that query."
                    )
                )
                return

            winning_file_name, encoded_question_tensor = ranking_result

            if winning_file_name:
                target_text_payload = AnswerSystem(
                    winning_file_name, encoded_question_tensor, query_text
                )
                is_loading = False
                self.status_box.after(
                    0, lambda: self.execute_typewriter_stream(
                        target_text_payload, winning_file_name
                    )
                )
            else:
                is_loading = False
                self.status_box.after(
                    0, lambda: self.render_fallback_msg(
                        "No matching notes found for that query."
                    )
                )

        threading.Thread(target=async_search_pipeline, daemon=True).start()

    def render_fallback_msg(self, msg):
        self.status_box.configure(state="normal")
        self.status_box.delete("1.0", ctk.END)
        self.status_box._textbox.insert("1.0", msg, "body")
        self.status_box.configure(state="disabled")
        self.search_btn.configure(state="normal")
        self.search_entry.configure(state="normal")

    def execute_typewriter_stream(self, target_text_payload, file_sources):
        self.status_box.configure(state="normal")
        self.status_box.delete("1.0", ctk.END)
        self.status_box.configure(state="disabled")

        payload_lines = target_text_payload.strip().splitlines()
        tb = self.status_box._textbox

        def clean(line):
            """Strip paired markdown bold/italic markers, leave lone chars."""
            s = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
            s = re.sub(r"__(.*?)__", s, s)
            s = re.sub(r"\*(.*?)\*", r"\1", s)
            return s

        def stream_line_by_line(line_idx=0):
            if line_idx < len(payload_lines):
                self.status_box.configure(state="normal")
                raw = payload_lines[line_idx]
                stripped = raw.strip()

                if stripped.startswith("## "):
                    # ── Section heading
                    text = clean(stripped[3:]).strip()
                    tb.insert(ctk.END, f"\n{text}\n", "h1")

                elif stripped.startswith("# "):
                    # ── Sub-heading
                    text = clean(stripped[2:]).strip()
                    tb.insert(ctk.END, f"\n{text}\n", "h2")

                elif stripped == "---":
                    # ── Thematic break — wide em-dash rule, not hyphens
                    tb.insert(ctk.END, "\n" + "─" * 56 + "\n", "divider")

                elif stripped.startswith("- "):
                    # ── Bullet — marker and text as separate tagged runs on one line
                    text = clean(stripped[2:]).strip()
                    tb.insert(ctk.END, "  ◆  ", "bullet_marker")
                    tb.insert(ctk.END, f"{text}\n", "bullet_text")

                elif stripped == "":
                    # ── Blank line — insert real breathing room
                    tb.insert(ctk.END, "\n")

                else:
                    tb.insert(ctk.END, f"{clean(raw)}\n", "body")

                self.status_box.see(ctk.END)
                self.status_box.configure(state="disabled")
                self.after(40, lambda: stream_line_by_line(line_idx + 1))

            else:
                # ── Source footer
                self.status_box.configure(state="normal")
                tb.insert(ctk.END, "\n" + "─" * 56 + "\n", "divider")
                tb.insert(ctk.END, "SOURCES\n", "source_label")

                sources = file_sources if isinstance(file_sources, list) else [file_sources]
                for src in sources:
                    tb.insert(ctk.END, f"  {src.upper()}\n", "source_item")

                self.status_box.see(ctk.END)
                self.status_box.configure(state="disabled")
                self.search_btn.configure(state="normal")
                self.search_entry.configure(state="normal")
                self.search_entry.focus()

        stream_line_by_line()


if __name__ == "__main__":
    app = StudyAssistantUI()
    app.mainloop()