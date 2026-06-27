"""Note loading, embedding, and caching pipeline.

This module is responsible for three things:

1. File ingestion — reads .txt, .md, .pdf, .docx, and .pptx files from
   the user's chosen notes folder and stores their plain-text content in
   the Notes dict.
2. Semantic embedding — encodes each note with a sentence-transformer model
   so similarity search can be performed without a round-trip to the LLM.
3. Disk cache — persists embeddings between sessions using pickle so the
   model doesn't re-encode unchanged files on every launch. The cache key
   is an MD5 hash of the file bytes so stale entries are detected automatically.

Module-level globals are used intentionally: loader is treated as a
shared-state namespace by retrieval.py and UI.py (they import loader and
read loader.Notes, loader.model, etc. directly).
"""

# loaders.py
import os
import json
import hashlib
import pickle
from pathlib import Path
from tkinter import filedialog
import pdfplumber
from docx import Document
from pptx import Presentation
from config import SETTINGS_FILE, EMBEDDINGS_CACHE_FILE, ALLOWED_EXTENSIONS
from network import check_internet_connection

# Shared state — imported by retrieval.py and UI.py
Notes = {}              # Filename → plain-text content for every loaded note
Embedding_cache = {}    # Filename → sentence-transformer tensor (in-memory)
Notes_Cache = {}        # Reserved for future use (currently unused)
model = None            # SentenceTransformer instance, set in bg_model_loading()
model_loaded = False    # True once bg_model_loading() finishes


def load_notes_from_path(folder_path):
    """Scan folder_path and load all supported note files into Notes.

    For each file the text content is extracted via the appropriate loader,
    the MD5 hash is compared against the on-disk embedding cache, and only
    changed files are re-encoded. Empty files are silently skipped; unreadable
    files are logged and skipped without raising.

    Args:
        folder_path: Absolute path to the folder containing the note files.
    """
    global Notes
    # Clear state so a folder switch doesn't mix old and new notes
    Notes.clear()
    Embedding_cache.clear()
    Notes_Cache.clear()
    disk_cache = {}

    # Load the on-disk embedding cache so we can skip re-encoding unchanged files
    if EMBEDDINGS_CACHE_FILE.exists():
        with open(EMBEDDINGS_CACHE_FILE, "rb") as f:
            disk_cache = pickle.load(f)
    notes_changed = False  # Track whether any file is new or modified

    for files in os.listdir(folder_path):
        path = os.path.join(folder_path, files)

        # Skip files with unsupported extensions
        if Path(path).suffix.lower() not in ALLOWED_EXTENSIONS:
            continue

        if os.path.isfile(path):
            try:
                # Dispatch to the appropriate format-specific loader
                if Path(path).suffix.lower() == ".pdf":
                    content = pdf_loader(path)
                elif Path(path).suffix.lower() == ".docx":
                    content = docx_loader(path)
                elif Path(path).suffix.lower() == ".pptx":
                    content = presentation_loader(path)
                else:
                    # Plain text or Markdown — read directly
                    with open(path, "r", encoding="utf-8") as f:
                        content = f.read()

                if not content:
                    print(f"Skipping empty file: {files}")
                    continue

                Notes[files] = content
                file_hash = get_hash(path)

                if files in disk_cache and disk_cache[files][0] == file_hash:
                    # File unchanged — reuse the cached embedding tensor
                    Embedding_cache[files] = disk_cache[files][1]
                else:
                    # File is new or modified — encode and update cache
                    Embedding_cache[files] = model.encode(
                        content, convert_to_tensor=True
                    )
                    disk_cache[files] = (file_hash, Embedding_cache[files])
                    notes_changed = True

            except Exception as e:
                print(f"Skipping unreadable asset {files}: {e}")

    # Evict cache entries for files that no longer exist in the folder
    disk_cache = {k: v for k, v in disk_cache.items() if k in Notes}

    # Only write to disk if something changed to avoid unnecessary I/O
    if notes_changed:
        with open(EMBEDDINGS_CACHE_FILE, "wb") as f:
            pickle.dump(disk_cache, f)


def load_folder_path():
    """Return the notes folder path, prompting the user if not yet configured.

    Checks SETTINGS_FILE for a previously saved path. If the saved path still
    exists on disk it is returned immediately. Otherwise an OS folder picker
    dialog is shown and the selection is persisted for future runs.

    Returns:
        The absolute path to the notes folder as a string, or None if the
        user cancels the dialog.
    """
    if SETTINGS_FILE.exists():
        with open(SETTINGS_FILE, "r", encoding="utf-8") as f:
            settings = json.load(f)
        path_notes = settings.get("path_notes", None)
        # Only return the saved path if it still exists on the filesystem
        if path_notes and os.path.exists(path_notes):
            return path_notes

    # No valid saved path — ask the user to choose a folder
    path_notes = filedialog.askdirectory(title="CHOOSE YOUR NOTES FOLDER.")
    if path_notes:
        with open(SETTINGS_FILE, "w", encoding="utf-8") as f:
            json.dump({"path_notes": path_notes}, f, indent=4)
        return path_notes
    return None


def bg_model_loading(ui_instance):
    """Load the embedding model and notes in a background thread.

    Intended to run in a threading.Thread so the splash screen remains
    responsive during the potentially slow model download and note-indexing
    steps. All widget updates use ui_instance.after() which is thread-safe.

    Args:
        ui_instance: The ACEUI instance whose status label should be updated.

    Returns:
        A (model, model_loaded) tuple; callers typically ignore the return
        value since both are accessible as module-level globals.
    """
    global model, model_loaded, is_online

    is_online = check_internet_connection()

    # Update splash screen — must use .after() because this runs off the main thread
    ui_instance.after(
        0, lambda: ui_instance.update_loading_status("LOADING EMBEDDINGS ENGINE...")
    )

    # Deferred import keeps cold-start time low when the model is already cached
    from sentence_transformers import SentenceTransformer

    model = SentenceTransformer("all-MiniLM-L6-v2")

    # Load the folder path and index notes (re-encodes only changed files)
    if load_folder_path():
        load_notes_from_path(load_folder_path())

    model_loaded = True

    # Signal the main thread to swap splash screen for main UI
    ui_instance.after(0, ui_instance.transition_to_main_ui)
    return model, model_loaded


def get_hash(filepath):
    """Compute the MD5 hex digest of a file's raw bytes.

    Used to detect whether a note file has changed since it was last encoded
    so the embedding cache can be invalidated selectively.

    Args:
        filepath: Absolute path to the file to hash.

    Returns:
        A lowercase hex string of the MD5 digest (32 characters).
    """
    with open(filepath, "rb") as f:
        return hashlib.md5(f.read()).hexdigest()


def pdf_loader(file_path):
    """Extract all plain text from a PDF file using pdfplumber.

    Iterates over every page and concatenates the extracted text. Pages that
    yield None (e.g. image-only pages) are silently skipped via the ``or ""``
    fallback.

    Args:
        file_path: Absolute path to the .pdf file.

    Returns:
        Concatenated text of all pages, stripped of leading/trailing whitespace.
    """
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text.strip()


def docx_loader(filepath):
    """Extract plain text from a .docx Word document.

    Reads paragraph objects from python-docx and concatenates their text.
    Formatting, images, and tables are not extracted.

    Args:
        filepath: Absolute path to the .docx file.

    Returns:
        Concatenated paragraph text, stripped of leading/trailing whitespace.
    """
    text = ""
    doc = Document(filepath)
    for paragraph in doc.paragraphs:
        text += paragraph.text or ""
    return text.strip()


def presentation_loader(file_path):
    """Extract plain text from a .pptx PowerPoint presentation.

    Iterates over all slides, shapes, and text-frame paragraphs. Only shapes
    that expose a text_frame attribute are read; images and charts are skipped.

    Args:
        file_path: Absolute path to the .pptx file.

    Returns:
        Concatenated slide text, stripped of leading/trailing whitespace.
    """
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text or ""
    return text.strip()
